import streamlit as st
import os
from pptx import Presentation
from transformers import MarianMTModel, AutoTokenizer
from io import BytesIO
from ArabicOcr import arabicocr
import pytesseract
import shutil
import os
import random

import pytesseract
from PIL import ImageEnhance, ImageFilter, Image


from pytesseract import Output
from PIL import Image
import cv2
from pdf2image import convert_from_bytes

def ocrcore(img, language):
    text= pytesseract.image_to_string(img, lang= language)
    return text

def get_grayscale(image):
    return cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

def remove_noise(image):
    return cv2.medianBlur(image, 5)

def threshholding(image):
    return cv2.threshold(image, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]

def perform_pdfocr(pdf_file):
    # Convert PDF to images
    images = convert_from_bytes(pdf_file.read())

    # Perform OCR on each image
    extracted_text = ""
    for image in images:
        text = pytesseract.image_to_string(image)
        extracted_text += text + "\n"

    return extracted_text

st.set_page_config(page_title='Total Translator', layout='wide', initial_sidebar_state='expanded')
#image1= image.open('./icons/IDSG.jpeg')
#image2= image.open('./icons/AI2Clogo.jpeg')
col1, col2, col3, col4, col5, col6, col7, col8, col9, col10 = st.columns(10)
with col1:
    st.image('./icons/IDSG.jpeg', width=140)
with col10:
    st.image('./icons/AI2Clogo.jpeg', width=140)
#st.image('./icons/IDSG.jpeg', width=140)
#st.image('./icons/AI2Clogo.jpeg', width=140,)
with col5:    
    st.title('Total Translator')
st.write('More translation options can be added upon request.')

uploaded_file = st.file_uploader('upload your powerpoint file here')



if uploaded_file:
    if '.pptx' in uploaded_file.name:
        filename = st.write('Your uploaded file is ready to translate')
        out_name = uploaded_file.name.replace('.pptx', '')
    elif '.jpg' in uploaded_file.name:
        file_contents = uploaded_file.getvalue()
        #img =cv2.imread('ukrain.jpg')
        
        # Choose a path to save the file
        temp_file_path = "temp_file.jpg"
        img = cv2.imread(temp_file_path)
        # Write the contents of the BytesIO object to a file
        with open(temp_file_path, "wb") as f:
            f.write(file_contents)

        st.write("File uploaded successfully!")
        st.write("File path:", temp_file_path)
    elif '.png' in uploaded_file.name:
        file_contents = uploaded_file.getvalue()
        #img =cv2.imread('ukrain.jpg')
        
        # Choose a path to save the file
        temp_file_path = "temp_file.png"
        img = cv2.imread(temp_file_path)
        # Write the contents of the BytesIO object to a file
        with open(temp_file_path, "wb") as f:
            f.write(file_contents)

        st.write("File uploaded successfully!")
        st.write("File path:", temp_file_path)
    elif '.pdf' in uploaded_file.name:


        st.write("PDF uploaded successfully!")
        
        
    
    else:
        st.error('Please upload a Powerpoint file ending in .pptx or .jpg')

Languages = {'arabic':'ar','english':'en', 'chinese': 'zh', 'ukrainian': 'uk', 'russian':'ru', 'placeholder':'none'}
languageocr = {'arabic':'ara','english':'eng', 'chinese': 'chi', 'ukrainian': 'ukr', 'russian':'rus', 'placeholder':'none'}
option1 = st.selectbox('Input language',
                      ('arabic', 'chinese', 'ukrainian','russian'))
option2 = st.selectbox('Output language',
                       ('english', 'placeholder'))

value1 = Languages[option1]
value2 = Languages[option2]

language= languageocr[option1]

    # Instantiate translation pipeline
def translation_pipeline(original_text):
    model_name = f"./opus-mt-{value1}-{value2}" #f"Helsinki-NLP/opus-mt-{value1}-{value2}"
    model = MarianMTModel.from_pretrained(model_name)
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    batch = tokenizer([original_text], return_tensors= 'pt')
    generated_ids = model.generate(**batch)
    translated_text = tokenizer.batch_decode(generated_ids, skip_special_tokens=True)[0]
    return translated_text

on = st.toggle('Chat Box for Direct Translation')

if on:
    user_input = st.text_input("Enter your message:", "")
    if st.button("Send"):
        results = translation_pipeline(user_input)
        st.write(results)

if st.button('Translate File'):
    
    # Instantiate translation pipeline
    def translation_pipeline(original_text):
        model_name = f"./opus-mt-{value1}-{value2}" #f"Helsinki-NLP/opus-mt-{value1}-{value2}"
        model = MarianMTModel.from_pretrained(model_name)
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        batch = tokenizer([original_text], return_tensors= 'pt')
        generated_ids = model.generate(**batch)
        translated_text = tokenizer.batch_decode(generated_ids, skip_special_tokens=True)[0]
        return translated_text

    # Load the presentation
    

    # Helsinki model has a bug that replaces a empty prompt with this obnoxious string.
    filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'

    if '.pptx' in uploaded_file.name:
        prs = Presentation(uploaded_file)
        # For each slide in the presentation
        for slide_number, slide in enumerate(prs.slides):  
            
            # For each shape in a slide
            for shape in slide.shapes:
                
                # Testing the "has_text_frame" parameter
                if shape.has_text_frame:

                    # For each paragraph of the text_frame
                    for paragraph in shape.text_frame.paragraphs:
                        
                        # Send paragraphs through translation pipeline
                        results = translation_pipeline(paragraph.text)
                        text_fixed = results.replace(filler, '')
                        paragraph.text = text_fixed
                
                elif shape.has_table:
                    table_filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'
                    tbl = shape.table
                    row_count = len(tbl.rows)
                    col_count = len(tbl.columns)
                    for r in range(0, row_count):
                        for c in range(0, col_count):
                            cell = tbl.cell(r, c).text_frame.fit_text(font_family='Arial', max_size=14, bold=False, italic=False)
                            paragraphs = cell.text_frame.paragraphs
                            for paragraph in paragraphs:
                                results = translation_pipeline(paragraph.text)
                                text_fixed = results.replace(filler, '')
                                paragraph.text = text_fixed
        binary_output = BytesIO()
        # save new file with translations
        prs.save(binary_output)

        st.success('Your Powerpoint file has been translated')
        st.download_button(label='Click to download PowerPoint',data=binary_output.getvalue(),file_name=f'{out_name}-translated.pptx')
    
    elif '.jpg' in uploaded_file.name:
        if value1 == 'ar':
            image_path= temp_file_path
            out_image='out.jpg'
            results=arabicocr.arabic_ocr(image_path,out_image)
            words=[]
            for i in range(len(results)):	
                    word=results[i][1]
                    words.append(word)
            translated= []
            translated= ' '.join(words)
            results = translation_pipeline(translated)
            text_to_add = results

            st.success('Your Image file has been translated')
            st.download_button('The Text from your image', text_to_add, file_name=f'Your_translated.txt')

        else:
            img = get_grayscale(img)
            img = threshholding(img)
            img = remove_noise(img)

            text = ocrcore(img, language)
            results = translation_pipeline(text)
            text_to_add = results

            st.success('Your Image file has been translated')
            st.download_button('The Text from your image', text_to_add, file_name=f'Your_translated.txt')

    elif '.png' in uploaded_file.name:
        if value1 == 'ar':
            image_path= temp_file_path
            out_image='out.png'
            results=arabicocr.arabic_ocr(image_path,out_image)
            words=[]
            for i in range(len(results)):	
                    word=results[i][1]
                    words.append(word)
            translated= []
            translated= ' '.join(words)
            results = translation_pipeline(translated)
            text_to_add = results

            st.success('Your Image file has been translated')
            st.download_button('The Text from your image', text_to_add, file_name=f'Your_translated.txt')

        else:
            img = get_grayscale(img)
            img = threshholding(img)
            img = remove_noise(img)

            text = ocrcore(img, language)
            results = translation_pipeline(text)
            text_to_add = results

            st.success('Your Image file has been translated')
            st.download_button('The Text from your image', text_to_add, file_name=f'Your_translated.txt')

    elif '.pdf' in uploaded_file.name:
        text = perform_pdfocr(uploaded_file)
        #st.download_button('temporary test', text, file_name='temporary_test.txt')
        results = translation_pipeline(text)
        text_to_add = results

        st.success('Your Image file has been translated')
        st.download_button('The Text from your image', text_to_add, file_name=f'Your_translated.txt')


    else:
        st.error("Unsupported File Type")
